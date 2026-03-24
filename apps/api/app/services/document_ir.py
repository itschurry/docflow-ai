from __future__ import annotations

import csv
from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation


def build_word_ir_from_markdown(title: str, body_text: str) -> dict[str, Any]:
    sections: list[dict[str, Any]] = []
    current = {"heading": title or "문서", "level": 1, "blocks": []}

    def flush() -> None:
        nonlocal current
        if current["blocks"]:
            sections.append(current)
        current = {"heading": title or "문서", "level": 1, "blocks": []}

    for raw in str(body_text or "").splitlines():
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("### "):
            flush()
            current = {"heading": stripped[4:].strip() or "세부 내용", "level": 3, "blocks": []}
            continue
        if stripped.startswith("## "):
            flush()
            current = {"heading": stripped[3:].strip() or "섹션", "level": 2, "blocks": []}
            continue
        if stripped.startswith("# "):
            flush()
            current = {"heading": stripped[2:].strip() or title or "문서", "level": 1, "blocks": []}
            continue
        if stripped.startswith(("- ", "* ")):
            current["blocks"].append({"type": "bullet", "text": stripped[2:].strip()})
        else:
            current["blocks"].append({"type": "paragraph", "text": stripped})
    if current["blocks"]:
        sections.append(current)

    return {
        "document_type": "word",
        "title": title or "문서",
        "sections": sections,
        "tables": [],
        "sources": extract_sources_from_text(body_text),
        "notes": [],
        "metadata": {"source": "markdown"},
    }


def build_sheet_ir_from_outline(title: str, rows: list[list[Any]], *, sheet_name: str = "summary") -> dict[str, Any]:
    normalized_rows = [[str(cell or "").strip() for cell in row] for row in rows if any(str(cell or "").strip() for cell in row)]
    return {
        "document_type": "sheet",
        "title": title or "시트 문서",
        "sheets": [
            {
                "name": sheet_name[:31] or "summary",
                "rows": normalized_rows,
            }
        ],
        "tables": [],
        "sources": [],
        "notes": [],
        "metadata": {"source": "outline"},
    }


def build_slides_ir(title: str, slide_outline: list[dict[str, Any]], *, sources: list[str] | None = None) -> dict[str, Any]:
    slides = []
    for index, slide in enumerate(slide_outline, start=1):
        slides.append(
            {
                "index": index,
                "title": str(slide.get("title") or f"슬라이드 {index}").strip(),
                "bullets": [str(item).strip() for item in slide.get("bullets") or [] if str(item).strip()],
                "speaker_notes": str(slide.get("speaker_notes") or "").strip(),
            }
        )
    return {
        "document_type": "slides",
        "title": title or "발표자료",
        "slides": slides,
        "sources": [str(item).strip() for item in (sources or []) if str(item).strip()],
        "notes": [],
        "metadata": {"source": "structured_deliverable"},
    }


def parse_document_to_ir(file_path: str, content_type: str | None = None) -> dict[str, Any]:
    suffix = Path(file_path).suffix.lower()
    title = Path(file_path).stem

    if suffix in {".txt", ".md", ".json"}:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        ir = build_word_ir_from_markdown(title, text)
        ir["metadata"]["content_type"] = content_type or ""
        return ir

    if suffix == ".csv":
        with Path(file_path).open("r", encoding="utf-8", errors="ignore", newline="") as fh:
            reader = csv.reader(fh)
            rows = [row for row in reader]
        return {
            "document_type": "sheet",
            "title": title,
            "sheets": [{"name": title[:31] or "sheet1", "rows": rows[:200]}],
            "tables": [],
            "sources": [],
            "notes": [],
            "metadata": {"source": "csv", "content_type": content_type or ""},
        }

    if suffix == ".docx":
        return _parse_docx_to_ir(file_path, content_type)

    if suffix == ".xlsx":
        return _parse_xlsx_to_ir(file_path, content_type)

    if suffix == ".pptx":
        return _parse_pptx_to_ir(file_path, content_type)

    if suffix == ".pdf":
        return _parse_pdf_to_ir(file_path, content_type)

    return {
        "document_type": "unknown",
        "title": title,
        "sections": [],
        "tables": [],
        "sources": [],
        "notes": [],
        "metadata": {"content_type": content_type or "", "suffix": suffix},
    }


def extract_text_from_ir(ir: dict[str, Any]) -> str:
    document_type = str(ir.get("document_type") or "")
    if document_type == "word":
        lines = [str(ir.get("title") or "").strip()]
        for section in ir.get("sections") or []:
            heading = str(section.get("heading") or "").strip()
            if heading:
                lines.append(heading)
            for block in section.get("blocks") or []:
                text = str(block.get("text") or "").strip()
                if text:
                    lines.append(text)
        for table in ir.get("tables") or []:
            for row in table.get("rows") or []:
                line = " | ".join(str(cell or "").strip() for cell in row if str(cell or "").strip())
                if line:
                    lines.append(line)
        return "\n".join(line for line in lines if line).strip()
    if document_type == "sheet":
        lines = [str(ir.get("title") or "").strip()]
        for sheet in ir.get("sheets") or []:
            name = str(sheet.get("name") or "").strip()
            if name:
                lines.append(f"[{name}]")
            for row in sheet.get("rows") or []:
                line = " | ".join(str(cell or "").strip() for cell in row if str(cell or "").strip())
                if line:
                    lines.append(line)
        return "\n".join(line for line in lines if line).strip()
    if document_type == "slides":
        lines = [str(ir.get("title") or "").strip()]
        for slide in ir.get("slides") or []:
            lines.append(str(slide.get("title") or "").strip())
            for bullet in slide.get("bullets") or []:
                if str(bullet).strip():
                    lines.append(str(bullet).strip())
            note = str(slide.get("speaker_notes") or "").strip()
            if note:
                lines.append(note)
        return "\n".join(line for line in lines if line).strip()
    return ""


def summarize_document_ir(ir: dict[str, Any]) -> str:
    document_type = str(ir.get("document_type") or "unknown")
    title = str(ir.get("title") or "문서").strip() or "문서"
    if document_type == "word":
        sections = ir.get("sections") or []
        tables = ir.get("tables") or []
        heading_preview = ", ".join(
            str(item.get("heading") or "").strip()
            for item in sections[:3]
            if str(item.get("heading") or "").strip()
        ) or "제목 미정"
        return f"{title} · word 문서 · 섹션 {len(sections)}개 · 표 {len(tables)}개 · 주요 섹션: {heading_preview}"
    if document_type == "sheet":
        sheets = ir.get("sheets") or []
        sheet_preview = ", ".join(str(item.get("name") or "").strip() for item in sheets[:3] if str(item.get("name") or "").strip()) or "시트 없음"
        return f"{title} · sheet 문서 · 시트 {len(sheets)}개 · 시트: {sheet_preview}"
    if document_type == "slides":
        slides = ir.get("slides") or []
        slide_preview = ", ".join(str(item.get("title") or "").strip() for item in slides[:3] if str(item.get("title") or "").strip()) or "슬라이드 없음"
        return f"{title} · slides 문서 · 슬라이드 {len(slides)}개 · 주요 슬라이드: {slide_preview}"
    return f"{title} · 구조 분석 불가"


def render_ir_to_docx_bytes(ir: dict[str, Any]) -> bytes:
    doc = Document()
    title = str(ir.get("title") or "Document").strip() or "Document"
    doc.add_heading(title, 0)
    document_type = str(ir.get("document_type") or "word")

    if document_type == "slides":
        for slide in ir.get("slides") or []:
            doc.add_heading(str(slide.get("title") or "Slide").strip(), level=1)
            for bullet in slide.get("bullets") or []:
                p = doc.add_paragraph(style="List Bullet")
                p.text = str(bullet).strip()
            note = str(slide.get("speaker_notes") or "").strip()
            if note:
                doc.add_paragraph(f"발표 포인트: {note}")
    elif document_type == "sheet":
        for sheet in ir.get("sheets") or []:
            doc.add_heading(str(sheet.get("name") or "Sheet").strip(), level=1)
            for row in sheet.get("rows") or []:
                values = [str(cell or "").strip() for cell in row if str(cell or "").strip()]
                if values:
                    doc.add_paragraph(" | ".join(values))
    else:
        for section in ir.get("sections") or []:
            doc.add_heading(str(section.get("heading") or "섹션").strip(), level=min(int(section.get("level") or 2), 3))
            for block in section.get("blocks") or []:
                text = str(block.get("text") or "").strip()
                if not text:
                    continue
                if block.get("type") == "bullet":
                    p = doc.add_paragraph(style="List Bullet")
                    p.text = text
                else:
                    doc.add_paragraph(text)
        for table in ir.get("tables") or []:
            rows = table.get("rows") or []
            if not rows:
                continue
            table_obj = doc.add_table(rows=len(rows), cols=max(len(row) for row in rows))
            for r_index, row in enumerate(rows):
                for c_index, cell in enumerate(row):
                    table_obj.cell(r_index, c_index).text = str(cell or "")

    sources = [str(item).strip() for item in ir.get("sources") or [] if str(item).strip()]
    if sources:
        doc.add_heading("참고 출처", level=1)
        for source in sources:
            p = doc.add_paragraph(style="List Bullet")
            p.text = source

    stream = BytesIO()
    doc.save(stream)
    return stream.getvalue()


def render_ir_to_xlsx_bytes(ir: dict[str, Any]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "overview"
    ws.append(["title", str(ir.get("title") or "Document")])
    ws.append(["document_type", str(ir.get("document_type") or "unknown")])
    ws.append([])

    document_type = str(ir.get("document_type") or "")
    if document_type == "sheet" and ir.get("sheets"):
        wb.remove(ws)
        for index, sheet in enumerate(ir.get("sheets") or []):
            ws_sheet = wb.active if index == 0 else wb.create_sheet()
            ws_sheet.title = (str(sheet.get("name") or f"sheet{index + 1}")[:31] or f"sheet{index + 1}")
            for row in sheet.get("rows") or []:
                ws_sheet.append([str(cell or "") for cell in row])
    elif document_type == "slides":
        ws.append(["slide", "bullet", "speaker_notes"])
        for slide in ir.get("slides") or []:
            title = str(slide.get("title") or "").strip()
            bullets = slide.get("bullets") or [""]
            note = str(slide.get("speaker_notes") or "").strip()
            for idx, bullet in enumerate(bullets):
                ws.append([title if idx == 0 else "", str(bullet or ""), note if idx == 0 else ""])
    else:
        ws.append(["section", "block_type", "text"])
        for section in ir.get("sections") or []:
            heading = str(section.get("heading") or "").strip()
            for idx, block in enumerate(section.get("blocks") or []):
                ws.append([heading if idx == 0 else "", str(block.get("type") or ""), str(block.get("text") or "")])
        if ir.get("tables"):
            table_sheet = wb.create_sheet("tables")
            for table_index, table in enumerate(ir.get("tables") or [], start=1):
                table_sheet.append([f"table_{table_index}"])
                for row in table.get("rows") or []:
                    table_sheet.append([str(cell or "") for cell in row])
                table_sheet.append([])

    sources = [str(item).strip() for item in ir.get("sources") or [] if str(item).strip()]
    if sources:
        src = wb.create_sheet("sources")
        src.append(["source"])
        for item in sources:
            src.append([item])

    stream = BytesIO()
    wb.save(stream)
    return stream.getvalue()


def render_ir_to_pptx_bytes(ir: dict[str, Any]) -> bytes:
    prs = Presentation()
    title = str(ir.get("title") or "Presentation").strip() or "Presentation"

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "DocFlow AI generated draft"

    slides = _coerce_ir_to_slide_outline(ir)
    bullet_layout = prs.slide_layouts[1]
    for item in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = item.get("title", "Untitled")
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        for index, line in enumerate(item.get("bullets", [])):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(line)
            paragraph.level = 0

    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def extract_sources_from_text(body_text: str | None) -> list[str]:
    sources: list[str] = []
    seen: set[str] = set()
    for raw in str(body_text or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        candidate = line[2:].strip() if line.startswith(("- ", "* ")) else line
        lowered = candidate.lower()
        if "http://" in lowered or "https://" in lowered or lowered.startswith("출처") or lowered.startswith("sources"):
            if lowered not in seen:
                seen.add(lowered)
                sources.append(candidate)
    return sources[:10]


def _parse_docx_to_ir(file_path: str, content_type: str | None) -> dict[str, Any]:
    doc = Document(file_path)
    title = Path(file_path).stem
    sections: list[dict[str, Any]] = []
    current = {"heading": title or "문서", "level": 1, "blocks": []}
    tables: list[dict[str, Any]] = []

    def flush() -> None:
        nonlocal current
        if current["blocks"]:
            sections.append(current)
        current = {"heading": title or "문서", "level": 1, "blocks": []}

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = str(paragraph.style.name or "").lower()
        if style_name.startswith("heading"):
            flush()
            level = 2
            if style_name[-1:].isdigit():
                level = int(style_name[-1:])
            current = {"heading": text, "level": min(level, 3), "blocks": []}
            continue
        block_type = "bullet" if "list" in style_name or paragraph.style.name == "List Bullet" else "paragraph"
        current["blocks"].append({"type": block_type, "text": text})
    if current["blocks"]:
        sections.append(current)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if rows:
            tables.append({"rows": rows})

    return {
        "document_type": "word",
        "title": title,
        "sections": sections,
        "tables": tables,
        "sources": extract_sources_from_text(extract_text_from_docx(doc)),
        "notes": [],
        "metadata": {"content_type": content_type or ""},
    }


def _parse_xlsx_to_ir(file_path: str, content_type: str | None) -> dict[str, Any]:
    wb = load_workbook(filename=file_path, data_only=True, read_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell) for cell in row]
            if any(str(cell).strip() for cell in values):
                rows.append(values[:20])
            if len(rows) >= 200:
                break
        sheets.append({"name": sheet.title[:31], "rows": rows})
    return {
        "document_type": "sheet",
        "title": Path(file_path).stem,
        "sheets": sheets,
        "tables": [],
        "sources": [],
        "notes": [],
        "metadata": {"content_type": content_type or ""},
    }


def _parse_pptx_to_ir(file_path: str, content_type: str | None) -> dict[str, Any]:
    prs = Presentation(file_path)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            lines = [line.strip() for line in shape.text.splitlines() if line.strip()]
            if not lines:
                continue
            if not title and getattr(shape, "is_placeholder", False) and getattr(shape.placeholder_format, "idx", None) == 0:
                title = lines[0]
                if len(lines) > 1:
                    bullets.extend(lines[1:])
            elif not title:
                title = lines[0]
                bullets.extend(lines[1:])
            else:
                bullets.extend(lines)
        slides.append(
            {
                "index": index,
                "title": title or f"슬라이드 {index}",
                "bullets": bullets[:8],
                "speaker_notes": "",
            }
        )
    return {
        "document_type": "slides",
        "title": Path(file_path).stem,
        "slides": slides,
        "sources": [],
        "notes": [],
        "metadata": {"content_type": content_type or ""},
    }


def _parse_pdf_to_ir(file_path: str, content_type: str | None) -> dict[str, Any]:
    try:
        import fitz
    except Exception:
        return {
            "document_type": "word",
            "title": Path(file_path).stem,
            "sections": [],
            "tables": [],
            "sources": [],
            "notes": [],
            "metadata": {"content_type": content_type or "", "parse_error": "fitz unavailable"},
        }
    doc = fitz.open(file_path)
    sections = []
    for index, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            sections.append(
                {
                    "heading": f"페이지 {index}",
                    "level": 2,
                    "blocks": [{"type": "paragraph", "text": line.strip()} for line in text.splitlines() if line.strip()][:40],
                }
            )
    doc.close()
    return {
        "document_type": "word",
        "title": Path(file_path).stem,
        "sections": sections,
        "tables": [],
        "sources": extract_sources_from_text(extract_text_from_ir({"document_type": "word", "title": Path(file_path).stem, "sections": sections, "tables": []})),
        "notes": [],
        "metadata": {"content_type": content_type or ""},
    }


def extract_text_from_docx(doc: Document) -> str:
    lines = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines).strip()


def _coerce_ir_to_slide_outline(ir: dict[str, Any]) -> list[dict[str, Any]]:
    document_type = str(ir.get("document_type") or "")
    if document_type == "slides":
        return [
            {
                "title": str(slide.get("title") or "슬라이드").strip(),
                "bullets": [str(item).strip() for item in slide.get("bullets") or [] if str(item).strip()],
            }
            for slide in ir.get("slides") or []
        ]
    if document_type == "sheet":
        slides = []
        for sheet in ir.get("sheets") or []:
            preview_rows = []
            for row in sheet.get("rows") or []:
                values = [str(cell or "").strip() for cell in row if str(cell or "").strip()]
                if values:
                    preview_rows.append(" | ".join(values))
                if len(preview_rows) >= 4:
                    break
            slides.append({"title": str(sheet.get("name") or "시트").strip(), "bullets": preview_rows or ["데이터 없음"]})
        return slides
    slides = []
    for section in ir.get("sections") or []:
        bullets = [str(block.get("text") or "").strip() for block in section.get("blocks") or [] if str(block.get("text") or "").strip()]
        slides.append({"title": str(section.get("heading") or "섹션").strip(), "bullets": bullets[:5] or ["내용 없음"]})
    return slides or [{"title": str(ir.get("title") or "문서").strip(), "bullets": ["내용 없음"]}]
