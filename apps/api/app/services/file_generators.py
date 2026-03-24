from io import BytesIO

from docx import Document
from docx.oxml.ns import qn
from openpyxl import Workbook
from pptx import Presentation

from app.services.document_ir import (
    _apply_korean_font,
    render_ir_to_docx_bytes,
    render_ir_to_pptx_bytes,
    render_ir_to_xlsx_bytes,
)


def generate_report_docx(title: str, body_text: str) -> bytes:
    """Convert markdown-formatted body_text to a styled DOCX document."""
    doc = Document()
    _apply_korean_font(doc)
    doc.add_heading(title, 0)

    for line in body_text.splitlines():
        stripped = line.rstrip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            p.text = stripped[2:]
        elif stripped.startswith("---"):
            doc.add_paragraph("─" * 40)
        elif stripped == "":
            # skip blank lines to avoid extra empty paragraphs
            continue
        else:
            doc.add_paragraph(stripped)

    stream = BytesIO()
    doc.save(stream)
    return stream.getvalue()


def generate_budget_xlsx(items: list[dict], total: int) -> bytes:
    wb = Workbook()
    ws_input = wb.active
    ws_input.title = "input"
    ws_input.append(["category", "name", "unit_cost",
                    "months", "rate", "amount"])

    for item in items:
        amount = int(item["unit_cost"] * item["months"] * item["rate"])
        ws_input.append([
            item.get("category", ""),
            item.get("name", ""),
            item.get("unit_cost", 0),
            item.get("months", 0),
            item.get("rate", 0),
            amount,
        ])

    ws_summary = wb.create_sheet("summary")
    ws_summary.append(["metric", "value"])
    ws_summary.append(["total", total])

    stream = BytesIO()
    wb.save(stream)
    return stream.getvalue()


def generate_pptx(title: str, slides: list[dict]) -> bytes:
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "DocFlow AI generated draft"

    bullet_layout = prs.slide_layouts[1]
    for item in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = item.get("title", "Untitled")
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        for line in item.get("bullets", []):
            p = body.add_paragraph()
            p.text = str(line)
            p.level = 0

    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def generate_structured_docx(document_ir: dict) -> bytes:
    return render_ir_to_docx_bytes(document_ir)


def generate_structured_xlsx(document_ir: dict) -> bytes:
    return render_ir_to_xlsx_bytes(document_ir)


def generate_structured_pptx(document_ir: dict) -> bytes:
    return render_ir_to_pptx_bytes(document_ir)
